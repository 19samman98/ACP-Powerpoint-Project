######################################################################
# Author: Sam McFarland
#
# Purpose: A program which takes photos and compiles them in a specialized PowerPoint for use with different donors.
######################################################################

from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os


def create_desktop_folder(folder_name):
    """
    Creates a desktop folder for the powerpoints to be stored at
    :param folder_name: The name of the folder you'd like to create
    :return:
    """
    # Get the path to the user's home directory
    path_to_user = os.path.expanduser("~")

    # Combine with the Desktop folder
    path_to_desktop = os.path.join(path_to_user, "Desktop")

    # Create the folder
    new_folder_path = os.path.join(path_to_desktop, folder_name)
    os.makedirs(new_folder_path, exist_ok=True)


def resize_image(image):
    """
    Resizes an image
    :param image: The path of the image to resize
    :return: None
    """

    # Open an image
    im = Image.open(image)

    # Specify the new size
    new_size = (1152, 648)

    # Resize the image
    resized_image = im.resize(new_size)

    # Save the resized image
    resized_image.save(image)


def construct_pp():
    """
    Constructs the not-personalized, base PowerPoint using photos that should be included in each PowerPoint.
    :return: None
    """
    # Creating a Presentation object
    ppt = Presentation()

    # Sets the dimensions for the Presentation
    ppt.slide_width = Inches(16)
    ppt.slide_height = Inches(9)

    # Goes through the photos and adds them to a PowerPoint
    for num in range(1, 16):
        # Giving Image path for each image
        img_path = 'Photos/Construction' + str(num) + '.jpg'

        # Selecting blank slide
        blank_slide_layout = ppt.slide_layouts[6]

        # Attaching slides to ppt
        slide = ppt.slides.add_slide(blank_slide_layout)

        # For no margins
        left = top = Inches(0)

        # adding images to each slide
        slide.shapes.add_picture(img_path, left, top)

    # save file
    ppt.save('base.pptx')


def add_donor(donor):
    """
    A function for adding specific donor photos to the base PowerPoint, saving them in a specific file.
    :param donor: The name of the donor photo to be added.
    :return:
    """
    # Sets the margins to 0
    left = top = Inches(0)

    # Sets the image path based on the donor parameter
    img_path = 'Personal_photos/' + donor + '.jpg'

    # Sets the presentation to the base created in construct_pp()
    ppt = Presentation("base.pptx")

    # Sets a slide layout and adds the donor photo to a new slide
    blank_slide_layout = ppt.slide_layouts[6]
    slide = ppt.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture(img_path, left, top)

    # Get the path to the user's home directory
    path_to_user = os.path.expanduser("~")

    # Combine with the Desktop folder
    path_to_desktop = os.path.join(path_to_user, "Desktop")

    # Saves the new specific donor PowerPoint to the folder on the desktop
    ppt.save(str(path_to_desktop) + "/Donor PowerPoints/" + donor + ".pptx")


def main():
    """
    A function that resizes the images to the correct size and creates personalized PowerPoints for donors.
    :return: None
    """
    # Makes the folder to put the PowerPoints in
    create_desktop_folder("Donor PowerPoints")

    # Resizes the images to the correct size
    for num in range(1, 16):
        resize_image('Photos/Construction' + str(num) + '.jpg')

    resize_image("Personal_photos/Jack.jpg")
    resize_image("Personal_photos/Joan.jpg")
    resize_image("Personal_photos/Koning.jpg")
    resize_image("Personal_photos/Miller.jpg")
    resize_image("Personal_photos/Shaw.jpg")
    resize_image("Personal_photos/Bateman.jpg")

    # Constructs the base PowerPoint and a PowerPoint for each donor
    construct_pp()
    add_donor("Jack")
    add_donor("Joan")
    add_donor("Koning")
    add_donor("Miller")
    add_donor("Shaw")
    add_donor("Bateman")

    print("Your PowerPoints have been created!")


if __name__ == "__main__":
    main()
